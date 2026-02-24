import fitz  # PyMuPDF
from pptx import Presentation
import pandas as pd
from typing import List, Dict, Any

def extract_pdf_text(file_path: str) -> List[Dict[str, Any]]:
    """Extracts text from PDF, returning page-level chunks with metadata."""
    doc = fitz.open(file_path)
    chunks = []
    
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        text = page.get_text("text").strip()
        if text:
            chunks.append({
                "text": text,
                "metadata": {"source": file_path, "page": page_num + 1}
            })
            
    return chunks

def extract_ppt_text(file_path: str) -> List[Dict[str, Any]]:
    """Extracts text from PowerPoint slides."""
    prs = Presentation(file_path)
    chunks = []
    
    for i, slide in enumerate(prs.slides):
        text_elements = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_elements.append(shape.text.strip())
                
        if text_elements:
            chunks.append({
                "text": "\n".join(text_elements),
                "metadata": {"source": file_path, "slide": i + 1}
            })
            
    return chunks

def extract_csv_data(file_path: str) -> List[Dict[str, Any]]:
    """Loads structured CSV data into dictionaries."""
    df = pd.read_csv(file_path)
    return df.to_dict(orient="records")